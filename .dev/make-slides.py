#!/usr/bin/env python3
"""Patch YC-20260511-cycif-tech-forum-ashlar.pptx with additional slides."""

from pathlib import Path
from pptx import Presentation
from pptx.util import Pt

SRC = Path("/Users/yuanchen/HMS Dropbox/Yu-An Chen/Slides/YC-20260511-cycif-tech-forum-ashlar.pptx")
DST = SRC.with_stem(SRC.stem + "-v2")

SIZE_L0 = Pt(24)   # 304800 EMU — matches existing level-0 bullets
SIZE_L1 = Pt(20)   # 254000 EMU — matches existing level-1 bullets


def get_layout(prs, name):
    return next(l for l in prs.slide_layouts if l.name == name)


def set_content(tf, items):
    """Overwrite a text frame. items: list of (text, level) tuples."""
    tf.clear()
    for i, (text, level) in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.level = level
        r = p.add_run()
        r.text = text
        r.font.size = SIZE_L0 if level == 0 else SIZE_L1


def move_slide(prs, from_idx, to_idx):
    lst = prs.slides._sldIdLst
    elems = list(lst)
    el = elems[from_idx]
    lst.remove(el)
    lst.insert(to_idx, el)


prs = Presentation(SRC)

# ── 1. Fix slide 2 title ────────────────────────────────────────────────────
# Original: "I have never run ASHLAR on my LSP desktop"
prs.slides[1].shapes.title.text = "Why people don’t run ASHLAR locally"

# ── 2. Rewrite slide 4 as the installation slide ────────────────────────────
# Original slide 4: "Facilitating local runs" (mixed pixi + GUI bullets)
# → becomes the "Getting started" installation slide; GUI features get their own slide
s4 = prs.slides[3]
s4.shapes.title.text = "Getting started"
set_content(s4.placeholders[1].text_frame, [
    ("Install pixi  →  pixi.sh", 0),
    ("One-time, no admin rights needed", 1),
    ("Set up the environment", 0),
    ("pixi install --locked", 1),
    ("Downloads Python, ashlar, and all dependencies into an isolated env", 1),
    ("Launch the GUI", 0),
    ("pixi run python run-ashlar.py", 1),
])

# ── 3. New slide: GUI features ───────────────────────────────────────────────
s = prs.slides.add_slide(get_layout(prs, "Title and Content"))
s.shapes.title.text = "What the GUI gives you"
set_content(s.placeholders[1].text_frame, [
    ("Samplesheet helper", 0),
    ("Auto-generates mcmicro CSV from a batch scan folder (LSP ID detection)", 1),
    ("Flexible input", 0),
    ("Directory layout or mcmicro samplesheet format", 1),
    ("Windows .lnk shortcut resolution", 1),
    ("Job management", 0),
    ("Parallel jobs, live log streaming, cancel mid-batch", 1),
    ("Output", 0),
    ("Pyramidal OME-TIFF + per-slide ashlar log written next to each slide", 1),
    ("Channel names embedded in OME-XML when a markers CSV is provided", 1),
])

# ── 4. New slide: GUI screenshot placeholder (Title Only) ────────────────────
s = prs.slides.add_slide(get_layout(prs, "Title Only"))
s.shapes.title.text = "The GUI"

# ── 5. New slide: Workflow overview ──────────────────────────────────────────
s = prs.slides.add_slide(get_layout(prs, "Title and Content"))
s.shapes.title.text = "Workflow overview"
set_content(s.placeholders[1].text_frame, [
    ("Input", 0),
    ("Scan folders (.rcpnl / .pysed.ome.tif) or mcmicro samplesheet CSV", 1),
    ("Samplesheet helper auto-generates the CSV from a batch scan folder", 1),
    ("Stitching", 0),
    ("run-ashlar GUI → one pyramidal OME-TIFF per sample", 1),
    ("Optional flat-field correction before stitching", 1),
    ("Channel names embedded in OME-XML after stitching", 1),
    ("Outputs feed directly into downstream mcmicro steps", 0),
    ("No extra file transfer — output written next to the scan folder", 1),
])

# ── Reorder: move Acknowledgement to the end ────────────────────────────────
# After the three adds, slide order is:
#   0 Title | 1 Pain points | 2 Use cases | 3 Installation
#   4 Acknowledgement (original) | 5 GUI features | 6 GUI demo | 7 Workflow
# → move Acknowledgement from index 4 to index 7 (last)
move_slide(prs, 4, 7)

prs.save(DST)
print(f"Saved: {DST}")
